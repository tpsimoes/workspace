"""Extract text from the three ASW CaaS Lead source files for skill authoring.
Outputs a single markdown-formatted digest to stdout.
"""
from __future__ import annotations
import sys
from pathlib import Path

from pptx import Presentation
from openpyxl import load_workbook

ROOT = Path(__file__).parent / "references"


def dump_pptx(path: Path) -> str:
    p = Presentation(str(path))
    lines = [f"# {path.name}", f"_slides: {len(p.slides)}_", ""]
    for i, slide in enumerate(p.slides, 1):
        lines.append(f"## Slide {i}")
        # collect text frames + tables
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(r.text for r in para.runs).strip()
                    if not t:
                        t = (para.text or "").strip()
                    if t:
                        lines.append(f"- {t}")
            if shape.has_table:
                tbl = shape.table
                lines.append("")
                lines.append("| " + " | ".join(cell.text.strip().replace("\n", " ")
                                                 for cell in tbl.rows[0].cells) + " |")
                lines.append("|" + "|".join("---" for _ in tbl.rows[0].cells) + "|")
                for row in list(tbl.rows)[1:]:
                    lines.append("| " + " | ".join(cell.text.strip().replace("\n", " ")
                                                     for cell in row.cells) + " |")
                lines.append("")
        # notes
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                lines.append("")
                lines.append(f"> **Notes:** {notes}")
        lines.append("")
    return "\n".join(lines)


def dump_xlsx(path: Path) -> str:
    wb = load_workbook(str(path), data_only=True, read_only=True)
    lines = [f"# {path.name}", f"_sheets: {wb.sheetnames}_", ""]
    for name in wb.sheetnames:
        ws = wb[name]
        lines.append(f"## Sheet: {name}")
        lines.append(f"_dims: {ws.max_row} rows x {ws.max_column} cols_")
        # Print first up to 30 rows
        for i, row in enumerate(ws.iter_rows(values_only=True), 1):
            if i > 30:
                lines.append(f"... ({ws.max_row - 30} more rows)")
                break
            vals = ["" if v is None else str(v)[:80] for v in row]
            lines.append("| " + " | ".join(vals) + " |")
        lines.append("")
    return "\n".join(lines)


def main() -> None:
    which = sys.argv[1] if len(sys.argv) > 1 else "all"
    files = {
        "stakeholder": ROOT / "ASWCustomerOutreach_StakeholderEngagementDeck.pptx",
        "charter":     ROOT / "ASWCustomerOutreach_FY26Dec_ProjectCharter.pptx",
        "kpis":        ROOT / "ASWCustomerOutreach_TargetCustomer_KPIs.xlsx",
    }
    parts = []
    for key, f in files.items():
        if which != "all" and which != key:
            continue
        if not f.exists():
            parts.append(f"MISSING: {f}")
            continue
        if f.suffix == ".pptx":
            parts.append(dump_pptx(f))
        elif f.suffix == ".xlsx":
            parts.append(dump_xlsx(f))
    print("\n\n---\n\n".join(parts))


if __name__ == "__main__":
    main()
